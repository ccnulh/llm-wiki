"""
素材导入模块
支持多种文档格式和音视频转文字
"""

import os
import re
import json
import time
import tempfile
import base64
import hashlib
from datetime import datetime
from pathlib import Path


# 支持的文件格式
SUPPORTED_FORMATS = {
    'document': ['.md', '.txt', '.json', '.csv', '.log'],
    'pdf': ['.pdf'],
    'word': ['.doc', '.docx'],
    'excel': ['.xls', '.xlsx'],
    'ppt': ['.ppt', '.pptx'],
    'ebook': ['.epub', '.mobi'],
    'audio': ['.mp3', '.wav', '.m4a', '.flac', '.aac', '.ogg', '.wma'],
    'video': ['.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.webm'],
    'image': ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'],
}


class Importer:
    """素材导入器"""

    def __init__(self, raw_dir: str, config: dict = None):
        self.raw_dir = raw_dir
        self.config = config or {}
        os.makedirs(raw_dir, exist_ok=True)

    def import_file(self, file_content: bytes, filename: str) -> dict:
        """智能导入文件，自动识别格式"""
        ext = os.path.splitext(filename)[1].lower()

        # 根据扩展名选择处理方式
        if ext in SUPPORTED_FORMATS['document']:
            return self._import_document(file_content, filename)
        elif ext in SUPPORTED_FORMATS['pdf']:
            return self.import_pdf(file_content, filename)
        elif ext in SUPPORTED_FORMATS['word']:
            return self.import_word(file_content, filename)
        elif ext in SUPPORTED_FORMATS['excel']:
            return self.import_excel(file_content, filename)
        elif ext in SUPPORTED_FORMATS['ppt']:
            return self.import_ppt(file_content, filename)
        elif ext in SUPPORTED_FORMATS['ebook']:
            return self.import_ebook(file_content, filename)
        elif ext in SUPPORTED_FORMATS['audio']:
            return self.import_audio(file_content, filename)
        elif ext in SUPPORTED_FORMATS['video']:
            return self.import_video(file_content, filename)
        elif ext in SUPPORTED_FORMATS['image']:
            return self.import_image(file_content, filename)
        else:
            # 未知格式，尝试作为文本处理
            return self._import_document(file_content, filename)

    def _import_document(self, file_content: bytes, filename: str) -> dict:
        """导入普通文档（md, txt, json, csv等）"""
        try:
            # 尝试解码为文本
            try:
                text = file_content.decode('utf-8')
            except:
                text = file_content.decode('gbk', errors='ignore')

            return self.import_text(text, title=filename, source=f'本地文件: {filename}')
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def import_text(self, content: str, title: str = None, source: str = None) -> dict:
        """导入文本内容"""
        timestamp = datetime.now().strftime('%Y-%m-%d_%H%M%S')
        safe_title = re.sub(r'[^\w\-]', '_', title or 'article')[:50]
        filename = f"{timestamp}_{safe_title}.md"

        md_content = f"""---
title: {title or '未知标题'}
source: {source or '未知来源'}
imported_at: {datetime.now().isoformat()}
---

{content}
"""
        filepath = os.path.join(self.raw_dir, filename)
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(md_content)

        return {
            'success': True,
            'filename': filename,
            'path': filepath
        }

    def import_pdf(self, file_content: bytes, filename: str) -> dict:
        """导入PDF文档"""
        timestamp = datetime.now().strftime('%Y-%m-%d_%H%M%S')
        safe_name = re.sub(r'[^\w\-\.]', '_', filename)

        # 保存原文
        pdf_path = os.path.join(self.raw_dir, f"{timestamp}_{safe_name}")
        with open(pdf_path, 'wb') as f:
            f.write(file_content)

        # 提取文字
        text = self._extract_pdf_text(pdf_path)

        return self.import_text(
            text,
            title=filename.replace('.pdf', ''),
            source=f'PDF文档: {filename}'
        )

    def _extract_pdf_text(self, pdf_path: str) -> str:
        """提取PDF文字"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n\n'
            return text.strip()
        except Exception as e:
            return f"[PDF文字提取失败: {str(e)}]"

    def import_word(self, file_content: bytes, filename: str) -> dict:
        """导入Word文档"""
        try:
            from docx import Document
            import io

            doc = Document(io.BytesIO(file_content))
            text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])

            return self.import_text(
                text,
                title=filename.replace('.docx', '').replace('.doc', ''),
                source=f'Word文档: {filename}'
            )
        except Exception as e:
            return {'success': False, 'error': f'Word解析失败: {str(e)}'}

    def import_excel(self, file_content: bytes, filename: str) -> dict:
        """导入Excel文档"""
        try:
            from openpyxl import load_workbook
            import io

            wb = load_workbook(io.BytesIO(file_content))
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"## {sheet_name}\n")

                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join([str(cell) if cell else '' for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)

            text = '\n'.join(text_parts)

            return self.import_text(
                text,
                title=filename.replace('.xlsx', '').replace('.xls', ''),
                source=f'Excel文档: {filename}'
            )
        except Exception as e:
            return {'success': False, 'error': f'Excel解析失败: {str(e)}'}

    def import_ppt(self, file_content: bytes, filename: str) -> dict:
        """导入PowerPoint文档"""
        try:
            from pptx import Presentation
            import io

            prs = Presentation(io.BytesIO(file_content))
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"## 幻灯片 {i}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            text = '\n'.join(text_parts)

            return self.import_text(
                text,
                title=filename.replace('.pptx', '').replace('.ppt', ''),
                source=f'PowerPoint文档: {filename}'
            )
        except Exception as e:
            return {'success': False, 'error': f'PPT解析失败: {str(e)}'}

    def import_ebook(self, file_content: bytes, filename: str) -> dict:
        """导入电子书（EPUB）"""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            import io

            ext = os.path.splitext(filename)[1].lower()

            if ext == '.epub':
                # 保存临时文件（epublib需要文件路径）
                with tempfile.NamedTemporaryFile(suffix='.epub', delete=False) as tmp:
                    tmp.write(file_content)
                    tmp_path = tmp.name

                try:
                    book = epub.read_epub(tmp_path)
                    text_parts = []

                    for item in book.get_items():
                        if item.get_type() == ebooklib.ITEM_DOCUMENT:
                            soup = BeautifulSoup(item.get_content(), 'html.parser')
                            text_parts.append(soup.get_text(separator='\n', strip=True))

                    text = '\n\n'.join(text_parts)
                finally:
                    os.unlink(tmp_path)

                return self.import_text(
                    text,
                    title=filename.replace('.epub', ''),
                    source=f'EPUB电子书: {filename}'
                )
            else:
                return {'success': False, 'error': f'不支持的电子书格式: {ext}'}

        except Exception as e:
            return {'success': False, 'error': f'电子书解析失败: {str(e)}'}

    def import_audio(self, file_content: bytes, filename: str) -> dict:
        """导入音频文件，使用阿里云语音识别转文字"""
        try:
            # 保存音频文件
            timestamp = datetime.now().strftime('%Y-%m-%d_%H%M%S')
            safe_name = re.sub(r'[^\w\-\.]', '_', filename)
            audio_path = os.path.join(self.raw_dir, f"{timestamp}_{safe_name}")

            with open(audio_path, 'wb') as f:
                f.write(file_content)

            # 调用阿里云语音识别
            transcript = self._speech_to_text(audio_path)

            return self.import_text(
                transcript,
                title=filename,
                source=f'音频转写: {filename}'
            )
        except Exception as e:
            return {'success': False, 'error': f'音频处理失败: {str(e)}'}

    def import_video(self, file_content: bytes, filename: str) -> dict:
        """导入视频文件，提取音频后转文字"""
        try:
            # 保存视频文件
            timestamp = datetime.now().strftime('%Y-%m-%d_%H%M%S')
            safe_name = re.sub(r'[^\w\-\.]', '_', filename)
            video_path = os.path.join(self.raw_dir, f"{timestamp}_{safe_name}")

            with open(video_path, 'wb') as f:
                f.write(file_content)

            # 提取音频
            audio_path = self._extract_audio_from_video(video_path)

            if not audio_path:
                return {'success': False, 'error': '无法从视频中提取音频'}

            # 调用阿里云语音识别
            transcript = self._speech_to_text(audio_path)

            # 删除临时音频文件
            if audio_path != video_path:
                try:
                    os.unlink(audio_path)
                except:
                    pass

            return self.import_text(
                transcript,
                title=filename,
                source=f'视频转写: {filename}'
            )
        except Exception as e:
            return {'success': False, 'error': f'视频处理失败: {str(e)}'}

    def import_image(self, file_content: bytes, filename: str) -> dict:
        """导入图片，使用OCR提取文字"""
        try:
            # 暂时只保存图片，OCR功能可选
            timestamp = datetime.now().strftime('%Y-%m-%d_%H%M%S')
            safe_name = re.sub(r'[^\w\-\.]', '_', filename)
            img_path = os.path.join(self.raw_dir, f"{timestamp}_{safe_name}")

            with open(img_path, 'wb') as f:
                f.write(file_content)

            # 尝试OCR（如果有配置）
            ocr_text = self._ocr_image(img_path)

            if ocr_text:
                return self.import_text(
                    ocr_text,
                    title=filename,
                    source=f'图片OCR: {filename}'
                )
            else:
                return {
                    'success': True,
                    'filename': os.path.basename(img_path),
                    'path': img_path,
                    'message': '图片已保存，OCR未配置或失败'
                }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def _extract_audio_from_video(self, video_path: str) -> str:
        """从视频中提取音频"""
        try:
            from moviepy.editor import VideoFileClip

            video = VideoFileClip(video_path)
            audio = video.audio

            if not audio:
                return None

            # 保存为wav（更好的兼容性）
            audio_path = video_path.rsplit('.', 1)[0] + '.wav'
            audio.write_audiofile(audio_path, verbose=False, logger=None, fps=16000)

            video.close()
            audio.close()

            return audio_path
        except Exception as e:
            print(f"提取音频失败: {e}")
            return None

    def _get_nls_token(self) -> str:
        """获取NLS Token"""
        asr_config = self.config.get('asr', {})

        # 优先使用缓存的Token
        token = asr_config.get('token', '')
        if token:
            return token

        # 使用AccessKey获取Token
        access_key_id = asr_config.get('access_key_id', '')
        access_key_secret = asr_config.get('access_key_secret', '')

        if access_key_id and access_key_secret:
            try:
                from token_getter import get_token_with_access_key
                result = get_token_with_access_key(access_key_id, access_key_secret)
                if result.get('success'):
                    token = result.get('token', '')
                    # 缓存Token
                    if token:
                        self._save_token(token)
                    return token
            except Exception as e:
                print(f"获取Token失败: {e}")

        return ''

    def _save_token(self, token: str):
        """保存Token到配置"""
        try:
            self.config['asr']['token'] = token
            config_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'config', 'settings.json')
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(self.config, f, indent=2, ensure_ascii=False)
        except:
            pass

    def _speech_to_text(self, audio_path: str) -> str:
        """将音频转为文字，优先使用本地Whisper，否则使用阿里云"""
        try:
            import requests
            import subprocess

            # 获取ASR配置
            asr_config = self.config.get('asr', {})

            # 优先使用本地Whisper
            if asr_config.get('use_local_whisper', False):
                return self._speech_to_text_whisper(audio_path)

            appkey = asr_config.get('appkey', '')

            if not appkey:
                return "[错误: 未配置阿里云语音识别AppKey]"

            # 获取Token
            token = self._get_nls_token()
            if not token:
                return "[错误: 无法获取语音识别Token，请检查AccessKey配置]"

            # 先将音频转换为16kHz PCM格式
            with tempfile.NamedTemporaryFile(suffix='.pcm', delete=False) as tmp:
                pcm_path = tmp.name

            try:
                # 使用ffmpeg转换为PCM格式
                cmd = [
                    'ffmpeg', '-y', '-i', audio_path,
                    '-f', 's16le',
                    '-acodec', 'pcm_s16le',
                    '-ar', '16000',
                    '-ac', '1',
                    pcm_path
                ]
                subprocess.run(cmd, capture_output=True, timeout=60)

                if not os.path.exists(pcm_path) or os.path.getsize(pcm_path) == 0:
                    return "[音频转换失败]"

                # 检查音频时长（PCM 16kHz 16bit 单声道 = 32000 bytes/sec）
                duration = os.path.getsize(pcm_path) / 32000
                if duration > 60:
                    # 长音频分段处理
                    return self._speech_to_text_long_pcm(pcm_path, token, appkey)

            except Exception as e:
                return f"[音频处理失败: {str(e)}]"

            # 使用阿里云一句话识别API
            url = "https://nls-gateway.cn-shanghai.aliyuncs.com/stream/v1/asr"

            headers = {
                "X-NLS-Token": token,
                "Content-Type": "audio/pcm;rate=16000",
            }

            params = {
                "appkey": appkey,
                "format": "pcm",
                "sample_rate": 16000,
                "enable_punctuation_prediction": "true",
                "enable_inverse_text_normalization": "true",
            }

            # 读取PCM数据
            with open(pcm_path, 'rb') as f:
                audio_data = f.read()

            # 清理临时文件
            try:
                os.unlink(pcm_path)
            except:
                pass

            # 发送请求
            response = requests.post(
                url,
                headers=headers,
                params=params,
                data=audio_data,
                timeout=120
            )

            if response.status_code == 200:
                result = response.json()
                if 'result' in result:
                    return result['result']
                else:
                    return f"[语音识别返回: {json.dumps(result, ensure_ascii=False)}]"
            else:
                return f"[语音识别失败: {response.status_code} - {response.text}]"

        except Exception as e:
            return f"[语音识别失败: {str(e)}]"

    def _speech_to_text_whisper(self, audio_path: str) -> str:
        """使用本地Whisper进行语音识别"""
        try:
            from faster_whisper import WhisperModel

            # 获取配置
            asr_config = self.config.get('asr', {})
            model_size = asr_config.get('whisper_model', 'medium')  # tiny, base, small, medium, large
            device = asr_config.get('whisper_device', 'cpu')  # cpu, cuda

            print(f"使用Whisper模型: {model_size}, 设备: {device}")

            # 根据设备选择
            compute_type = "int8" if device == "cpu" else "float16"

            # 加载模型
            model = WhisperModel(
                model_size,
                device=device,
                compute_type=compute_type
            )

            # 转写
            segments, info = model.transcribe(
                audio_path,
                language='zh',  # 优先中文
                beam_size=5,
                vad_filter=True,  # 语音活动检测，过滤静音
                vad_parameters=dict(min_silence_duration_ms=500)
            )

            print(f"检测到语言: {info.language}，时长: {info.duration:.2f}秒")

            # 收集结果
            results = []
            for segment in segments:
                text = segment.text.strip()
                if text:
                    # 添加时间戳
                    start = int(segment.start)
                    end = int(segment.end)
                    results.append(f"[{start//60}:{start%60:02d}-{end//60}:{end%60:02d}] {text}")

            transcript = '\n'.join(results)

            if not transcript:
                return "[Whisper未能识别出文本]"

            return transcript

        except Exception as e:
            return f"[Whisper转写失败: {str(e)}]"

    def _speech_to_text_long_pcm(self, pcm_path: str, token: str, appkey: str) -> str:
        """处理长音频PCM文件"""
        import requests

        # 读取PCM数据
        with open(pcm_path, 'rb') as f:
            audio_data = f.read()

        # 分段处理，每段约50秒
        segment_size = 50 * 32000  # 50秒的数据
        segments = [audio_data[i:i+segment_size] for i in range(0, len(audio_data), segment_size)]

        results = []
        url = "https://nls-gateway.cn-shanghai.aliyuncs.com/stream/v1/asr"

        headers = {
            "X-NLS-Token": api_key,
            "Content-Type": "audio/pcm;rate=16000",
        }

        params = {
            "appkey": appkey,
            "format": "pcm",
            "sample_rate": 16000,
            "enable_punctuation_prediction": "true",
            "enable_inverse_text_normalization": "true",
        }

        for i, segment in enumerate(segments):
            try:
                response = requests.post(url, headers=headers, params=params, data=segment, timeout=120)
                if response.status_code == 200:
                    result = response.json()
                    if 'result' in result and result['result']:
                        results.append(result['result'])
            except Exception as e:
                print(f"分段{i+1}识别失败: {e}")

        # 清理临时文件
        try:
            os.unlink(pcm_path)
        except:
            pass

        return '\n'.join(results) if results else "[语音识别: 无法识别内容]"

    def _speech_to_text_long(self, audio_path: str, api_key: str) -> str:
        """处理长音频文件"""
        try:
            import requests

            # 对于长音频，使用分段处理
            # 将音频分割成多个片段，分别识别

            # 首先尝试使用ffmpeg分割音频
            import subprocess

            # 获取音频时长
            probe_cmd = [
                'ffprobe', '-v', 'error',
                '-show_entries', 'format=duration',
                '-of', 'default=noprint_wrappers=1:nokey=1',
                audio_path
            ]

            try:
                duration = float(subprocess.check_output(probe_cmd).decode().strip())
            except:
                duration = 300  # 默认5分钟

            # 每60秒分割一次
            segment_duration = 60
            num_segments = int(duration / segment_duration) + 1

            results = []

            for i in range(num_segments):
                start_time = i * segment_duration

                # 创建临时片段
                with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp:
                    segment_path = tmp.name

                try:
                    # 使用ffmpeg分割
                    split_cmd = [
                        'ffmpeg', '-y', '-i', audio_path,
                        '-ss', str(start_time),
                        '-t', str(segment_duration),
                        '-ar', '16000',
                        '-ac', '1',
                        segment_path
                    ]
                    subprocess.run(split_cmd, capture_output=True, timeout=30)

                    if os.path.getsize(segment_path) > 0:
                        # 识别这个片段
                        result = self._speech_to_text_segment(segment_path, api_key)
                        if result and not result.startswith('['):
                            results.append(result)

                finally:
                    if os.path.exists(segment_path):
                        os.unlink(segment_path)

            return '\n\n'.join(results) if results else "[语音识别: 无法处理长音频]"

        except Exception as e:
            return f"[长音频处理失败: {str(e)}]"

    def _speech_to_text_segment(self, segment_path: str, api_key: str) -> str:
        """识别单个音频片段"""
        import requests

        url = "https://nls-gateway.cn-shanghai.aliyuncs.com/stream/v1/asr"

        headers = {
            "X-NLS-Token": api_key,
            "Content-Type": "audio/wav;codec=pcm;rate=16000",
        }

        params = {
            "appkey": "nls-service",
            "format": "wav",
            "sample_rate": 16000,
            "enable_punctuation_prediction": True,
            "enable_inverse_text_normalization": True,
        }

        with open(segment_path, 'rb') as f:
            audio_data = f.read()

        response = requests.post(url, headers=headers, params=params, data=audio_data, timeout=30)

        if response.status_code == 200:
            result = response.json()
            return result.get('result', '')
        else:
            return ''

    def _speech_to_text_dashscope(self, audio_path: str, api_key: str) -> str:
        """使用DashScope API进行语音识别"""
        try:
            import requests

            url = "https://dashscope.aliyuncs.com/api/v1/services/audio/asr/transcription"

            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            }

            # 使用文件内容base64编码
            with open(audio_path, 'rb') as f:
                audio_data = f.read()

            audio_base64 = base64.b64encode(audio_data).decode('utf-8')

            data = {
                "model": "paraformer-8k-v1",
                "input": {
                    "audio": audio_base64
                }
            }

            response = requests.post(url, headers=headers, json=data, timeout=120)

            if response.status_code == 200:
                result = response.json()
                if 'output' in result and 'results' in result['output']:
                    texts = []
                    for item in result['output']['results']:
                        if 'transcription_text' in item:
                            texts.append(item['transcription_text'])
                    return '\n'.join(texts)
                else:
                    return f"[DashScope返回: {result}]"
            else:
                return f"[DashScope调用失败: {response.status_code} - {response.text}]"

        except Exception as e:
            return f"[DashScope识别失败: {str(e)}]"

    def _ocr_image(self, img_path: str) -> str:
        """OCR图片提取文字（可选功能）"""
        # 暂不实现，需要额外配置
        return None


class WebFetcher:
    """网页抓取器"""

    def fetch_url(self, url: str) -> dict:
        """抓取普通网页"""
        try:
            import requests
            from bs4 import BeautifulSoup

            # 知乎需要特殊处理
            if 'zhihu.com' in url:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (iPhone; CPU iPhone OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Mobile/15E148 Safari/604.1',
                    'Accept': 'text/html,application/xhtml+xml',
                }
            else:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
                    'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                }
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()

            soup = BeautifulSoup(response.text, 'html.parser')

            # 提取标题
            title = soup.find('title')
            title_text = title.string if title else '未知标题'

            # 移除脚本和样式
            for script in soup(['script', 'style', 'nav', 'footer', 'header']):
                script.decompose()

            # 提取正文
            content = ''
            for selector in ['article', '.content', '.post', '.article', 'main', '#content', '.entry-content']:
                elem = soup.select_one(selector)
                if elem:
                    content = elem.get_text(separator='\n', strip=True)
                    break

            if not content:
                body = soup.find('body')
                if body:
                    content = body.get_text(separator='\n', strip=True)

            return {
                'success': True,
                'title': title_text,
                'content': content,
                'source': url
            }

        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }


class WechatFetcher:
    """微信公众号文章抓取器 - 使用HTTP请求"""

    def fetch_article(self, url: str) -> dict:
        """抓取微信公众号文章"""
        try:
            import requests
            from bs4 import BeautifulSoup

            # 首先尝试直接抓取
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                'Referer': 'https://mp.weixin.qq.com/',
            }

            response = requests.get(url, headers=headers, timeout=30, allow_redirects=True)
            response.encoding = 'utf-8'

            # 检查是否被重定向到验证码页面
            if 'wappoc_appmsgcaptcha' in response.url or 'web.archive.org' in response.url:
                # 尝试使用 Google Cache 或其他方式
                cached_url = f"https://webcache.googleusercontent.com/search?q=cache:{url}"
                try:
                    cached_response = requests.get(cached_url, headers=headers, timeout=30)
                    if cached_response.status_code == 200:
                        response = cached_response
                        response.encoding = 'utf-8'
                except:
                    pass

            soup = BeautifulSoup(response.text, 'html.parser')

            # 提取标题
            title = ''
            title_elem = soup.select_one('.rich_media_title') or soup.select_one('title')
            if title_elem:
                title = title_elem.get_text(strip=True)

            # 提取正文
            content = ''
            content_elem = soup.select_one('#js_content') or soup.select_one('.rich_media_content')
            if content_elem:
                # 移除脚本和样式
                for tag in content_elem.find_all(['script', 'style', 'iframe']):
                    tag.decompose()
                content = content_elem.get_text(separator='\n', strip=True)

            # 如果还是获取不到，尝试另一种方式提取
            if not content:
                # 尝试从meta description获取摘要
                meta_desc = soup.select_one('meta[name="description"]')
                if meta_desc:
                    content = meta_desc.get('content', '')
                    title = soup.select_one('title')
                    if title:
                        title = title.get_text(strip=True)

            if not content:
                # 最后尝试：提取所有p标签内容
                paragraphs = soup.find_all('p')
                if paragraphs:
                    content = '\n'.join([p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)])

            if not content:
                return {'success': False, 'error': '无法提取文章内容，可能需要登录或验证码验证'}

            # 清理内容
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            content = '\n'.join(lines)

            return {
                'success': True,
                'title': title or '未知标题',
                'content': content,
                'source': url
            }

        except Exception as e:
            return {'success': False, 'error': str(e)}


class VideoFetcher:
    """视频内容抓取器 - 支持 YouTube, Bilibili 等平台"""

    def parse_video_url(self, url: str) -> tuple:
        """解析视频链接，返回平台和标准URL"""
        import re

        # YouTube
        yt_match = re.search(r'(?:youtube\.com\/watch\?v=|youtu\.be\/)([a-zA-Z0-9_-]{11})', url)
        if yt_match:
            video_id = yt_match.group(1)
            return 'youtube', video_id, f'https://www.youtube.com/watch?v={video_id}'

        # Bilibili
        bl_match = re.search(r'bilibili\.com\/video\/(BV[a-zA-Z0-9]+|av\d+)', url)
        if bl_match:
            vid = bl_match.group(1)
            return 'bilibili', vid, f'https://www.bilibili.com/video/{vid}'

        return None, None, None

    def fetch_video_content(self, url: str) -> dict:
        """抓取视频内容（字幕转录）"""
        import subprocess
        import json
        import tempfile
        from pathlib import Path

        platform, video_id, clean_url = self.parse_video_url(url)

        if not platform:
            return {'success': False, 'error': '不支持的视频链接，仅支持 YouTube 和 Bilibili'}

        try:
            # 获取视频信息
            info_cmd = ['yt-dlp', '--dump-json', '--no-download', clean_url]
            result = subprocess.run(info_cmd, capture_output=True, text=True, timeout=60)

            if result.returncode != 0:
                return {'success': False, 'error': '获取视频信息失败'}

            video_info = json.loads(result.stdout)
            title = video_info.get('title', '未知标题')
            uploader = video_info.get('uploader', '未知作者')
            duration = video_info.get('duration_string', '未知时长')
            description = video_info.get('description', '')

            # 创建临时目录
            temp_dir = tempfile.mkdtemp()

            # 下载字幕
            sub_cmd = [
                'yt-dlp',
                '--skip-download',
                '--write-sub',
                '--write-auto-sub',
                '--sub-lang', 'zh-Hans,zh-Hant,en,zh,ja',
                '--sub-format', 'srt',
                '-o', f'{temp_dir}/%(title)s.%(ext)s',
                clean_url
            ]

            sub_result = subprocess.run(sub_cmd, capture_output=True, text=True, timeout=120)

            # 查找字幕文件
            srt_files = list(Path(temp_dir).glob('*.srt'))

            transcript = ''
            if srt_files:
                # 选择最新的字幕文件
                srt_file = max(srt_files, key=lambda f: f.stat().st_mtime)
                with open(srt_file, 'r', encoding='utf-8') as f:
                    srt_content = f.read()
                transcript = self._srt_to_text(srt_content)

                # 清理临时文件
                for f in srt_files:
                    f.unlink()
            else:
                # 如果没有字幕，使用描述或提示
                if description:
                    transcript = f"[视频描述]\n{description}\n\n[注意] 此视频无内置字幕，仅提取了视频描述"
                else:
                    transcript = "[注意] 此视频无内置字幕，无法获取详细内容"

            # 组装内容
            content = f"""# {title}

**作者**: {uploader}
**时长**: {duration}
**平台**: {platform}
**链接**: {clean_url}

---

## 视频讲稿

{transcript}

---
*来源: {platform} 视频*
"""

            return {
                'success': True,
                'title': title,
                'content': content,
                'source': clean_url,
                'platform': platform,
                'has_transcript': bool(srt_files)
            }

        except subprocess.TimeoutExpired:
            return {'success': False, 'error': '视频处理超时'}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def _srt_to_text(self, srt_content: str) -> str:
        """将 SRT 格式转换为纯文本"""
        lines = srt_content.split('\n')
        text_lines = []

        for line in lines:
            line = line.strip()
            # 跳过序号、时间轴和空行
            if not line or line.isdigit() or '-->' in line:
                continue
            text_lines.append(line)

        return '\n'.join(text_lines)


class PodcastFetcher:
    """播客内容抓取器 - 支持 RSS 订阅源"""

    def fetch_rss(self, rss_url: str) -> dict:
        """抓取播客RSS并下载单集音频"""
        import xml.etree.ElementTree as ET
        import requests
        import re
        import tempfile

        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'
            }
            response = requests.get(rss_url, headers=headers, timeout=30)
            response.raise_for_status()

            # 解析RSS
            root = ET.fromstring(response.content)

            # 获取播客信息
            channel = root.find('channel')
            podcast_title = channel.find('title').text if channel.find('title') is not None else '未知播客'

            # 如果没有传递具体集数，返回播客信息供用户选择
            items = channel.findall('item')
            episodes = []

            for item in items[:10]:  # 最多返回10集
                episode_title = item.find('title').text if item.find('title') is not None else '未知标题'
                episode_desc = ''
                if item.find('description') is not None:
                    episode_desc = item.find('description').text or ''
                elif item.find('{}description') is not None:
                    episode_desc = item.find('{}description').text or ''

                # 查找音频链接
                audio_url = None
                enclosure = item.find('enclosure')
                if enclosure is not None and enclosure.get('url'):
                    audio_url = enclosure.get('url')

                # 兼容命名空间
                if not audio_url:
                    for child in item:
                        if 'enclosure' in child.tag and child.get('url'):
                            audio_url = child.get('url')
                            break

                pub_date = item.find('pubDate').text if item.find('pubDate') is not None else ''

                episodes.append({
                    'title': episode_title,
                    'description': episode_desc[:200] if episode_desc else '',
                    'audio_url': audio_url,
                    'pub_date': pub_date
                })

            return {
                'success': True,
                'type': 'podcast',
                'platform': 'rss',
                'title': podcast_title,
                'episodes': episodes,
                'message': f'获取到播客 "{podcast_title}"，共 {len(episodes)} 集'
            }

        except Exception as e:
            return {'success': False, 'error': f'RSS抓取失败: {str(e)}'}

    def download_episode(self, audio_url: str, episode_title: str, podcast_title: str) -> dict:
        """下载播客单集音频并转写"""
        import requests
        import tempfile
        import os

        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'
            }

            print(f"下载播客音频: {audio_url}")

            # 下载音频
            response = requests.get(audio_url, headers=headers, timeout=120, stream=True)
            response.raise_for_status()

            # 获取文件名
            content_disposition = response.headers.get('content-disposition', '')
            filename = re.findall('filename="?(.+?)"?', content_disposition)
            if filename:
                filename = filename[0]
            else:
                ext = audio_url.split('?')[0].split('.')[-1]
                ext = ext if ext in ['mp3', 'm4a', 'wav', 'ogg'] else 'mp3'
                filename = f"{episode_title[:50]}.{ext}"

            # 保存到临时文件
            temp_dir = tempfile.mkdtemp()
            audio_path = os.path.join(temp_dir, filename)

            with open(audio_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            file_size = os.path.getsize(audio_path)
            if file_size < 1000:
                return {'success': False, 'error': '音频文件太小，可能是下载失败'}

            return {
                'success': True,
                'audio_path': audio_path,
                'title': episode_title,
                'podcast_title': podcast_title,
                'filename': filename,
                'file_size': file_size
            }

        except Exception as e:
            return {'success': False, 'error': f'下载失败: {str(e)}'}


class FeishuFetcher:
    """飞书文档抓取器"""

    def __init__(self):
        self.app_id = None
        self.app_secret = None

    def fetch(self, url: str) -> dict:
        """抓取飞书文档内容"""
        try:
            import lark_oapi as lark
            import os
            import re
            from lark_oapi.api.docx.v1.model import GetDocumentRequest, ListDocumentBlockRequest

            # 从配置读取app_id和app_secret
            config_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'config', 'settings.json')
            config = {}
            if os.path.exists(config_path):
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)

            feishu_config = config.get('feishu', {})
            self.app_id = feishu_config.get('app_id', '')
            self.app_secret = feishu_config.get('app_secret', '')

            if not self.app_id or not self.app_secret:
                return {'success': False, 'error': '飞书应用凭证未配置，请在设置页面配置 app_id 和 app_secret'}

            # 从URL提取文档token
            match = re.search(r'/docx/([a-zA-Z0-9]+)', url)
            if not match:
                return {'success': False, 'error': '无效的飞书文档链接'}

            document_id = match.group(1)

            # 初始化lark客户端
            cli = lark.Client.builder().app_id(self.app_id).app_secret(self.app_secret).build()

            # 获取文档元信息
            doc_request = GetDocumentRequest.builder().document_id(document_id).build()
            doc_response = cli.docx.v1.document.get(doc_request)

            title = '飞书文档'
            if doc_response.success() and doc_response.data and doc_response.data.document:
                title = getattr(doc_response.data.document, 'title', '飞书文档')

            # 获取文档所有块 (使用 list API 获取所有块)
            block_request = ListDocumentBlockRequest.builder().document_id(document_id).build()
            block_response = cli.docx.v1.document_block.list(block_request)

            all_blocks = []
            if block_response.success() and block_response.data and block_response.data.items:
                all_blocks.extend(block_response.data.items)

            content = self._parse_blocks(all_blocks)

            if not content.strip():
                return {'success': False, 'error': '文档内容为空或无法解析'}

            return {
                'success': True,
                'title': title,
                'content': content,
                'source': url
            }

        except ImportError:
            return {'success': False, 'error': '飞书SDK未安装，请运行: pip install lark-oapi'}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def _parse_blocks(self, blocks: list) -> str:
        """解析飞书文档块为markdown"""
        md_parts = []

        for block in blocks:
            text = ''

            # Block type 1 = Page (文档标题页)
            if block.block_type == 1:
                # 尝试从 page 属性获取标题
                if hasattr(block, 'page') and block.page and hasattr(block.page, 'elements'):
                    for elem in block.page.elements:
                        if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                            text += elem.text_run.content or ''
                # 文档标题使用 page block 的内容
                if text:
                    md_parts.append(f"# {text}")

            # Block type 2 = Text
            elif block.block_type == 2:
                if hasattr(block, 'text') and block.text and hasattr(block.text, 'elements'):
                    for elem in block.text.elements:
                        if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                            text += elem.text_run.content or ''
                if text:
                    md_parts.append(text)

            # Block type 30 = Sheet (电子表格嵌入)
            elif block.block_type == 30:
                if hasattr(block, 'sheet') and block.sheet:
                    sheet_token = getattr(block.sheet, 'token', None)
                    if sheet_token:
                        # 尝试获取表格内容
                        sheet_content = self._fetch_sheet_content(sheet_token)
                        if sheet_content:
                            md_parts.append(sheet_content)
                        else:
                            md_parts.append(f"[表格嵌入: {sheet_token}]")

            # 其他文本类型块 (heading, paragraph, todo, bullet, ordered等)
            else:
                # 尝试获取text属性
                if hasattr(block, 'text') and block.text and hasattr(block.text, 'elements'):
                    for elem in block.text.elements:
                        if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                            text += elem.text_run.content or ''

                # 检查 heading1-9
                if not text:
                    for i in range(1, 10):
                        heading = getattr(block, f'heading{i}', None)
                        if heading and hasattr(heading, 'elements'):
                            for elem in heading.elements:
                                if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                                    text += elem.text_run.content or ''
                            if text:
                                text = '#' * i + ' ' + text
                                break

                # 检查 paragraph
                if not text and hasattr(block, 'paragraph') and block.paragraph:
                    para = block.paragraph
                    if hasattr(para, 'elements') and para.elements:
                        for elem in para.elements:
                            if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                                text += elem.text_run.content or ''

                # 检查 todo
                if not text and hasattr(block, 'todo') and block.todo:
                    td = block.todo
                    if hasattr(td, 'elements') and td.elements:
                        for elem in td.elements:
                            if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                                text += '[x] ' if getattr(td, 'done_status', False) else '[ ] '
                                text += elem.text_run.content or ''

                # 检查 bullet
                if not text and hasattr(block, 'bullet') and block.bullet:
                    bd = block.bullet
                    if hasattr(bd, 'elements') and bd.elements:
                        for elem in bd.elements:
                            if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                                text += '- ' + (elem.text_run.content or '')
                                break

                # 检查 ordered
                if not text and hasattr(block, 'ordered') and block.ordered:
                    od = block.ordered
                    if hasattr(od, 'elements') and od.elements:
                        for elem in od.elements:
                            if hasattr(elem, 'text_run') and elem.text_run and hasattr(elem.text_run, 'content'):
                                text += '1. ' + (elem.text_run.content or '')
                                break

                if text:
                    md_parts.append(text)

        return '\n\n'.join(md_parts)

    def _fetch_sheet_content(self, sheet_token: str) -> str:
        """获取飞书表格内容"""
        try:
            import requests

            # 获取 access token
            auth_url = "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal"
            auth_data = {"app_id": self.app_id, "app_secret": self.app_secret}
            auth_resp = requests.post(auth_url, json=auth_data, timeout=10)
            if auth_resp.status_code != 200:
                return None
            auth_json = auth_resp.json()
            if auth_json.get('code') != 0:
                return None
            access_token = auth_json.get('tenant_access_token')

            headers = {"Authorization": f"Bearer {access_token}"}

            # 解析 spreadsheet token 和 sheet id
            parts = sheet_token.rsplit('_', 1)
            if len(parts) != 2:
                return None
            spreadsheet_token = parts[0]
            sheet_id = parts[1]

            # 获取 sheets 列表找到实际的 sheet id
            sheets_url = f"https://open.feishu.cn/open-apis/sheets/v3/spreadsheets/{spreadsheet_token}/sheets/query"
            sheets_resp = requests.get(sheets_url, headers=headers, timeout=10)
            if sheets_resp.status_code != 200:
                return None
            sheets_json = sheets_resp.json()
            if sheets_json.get('code') != 0:
                return None

            actual_sheet_id = None
            for sheet in sheets_json.get('data', {}).get('sheets', []):
                # 使用 Sheet1 或第一个 sheet
                title = sheet.get('title', '')
                sid = sheet.get('sheet_id', '')
                if title == 'Sheet1' or title == sheet_id:
                    actual_sheet_id = sid
                    break
            if not actual_sheet_id:
                # 使用第一个可用的 sheet
                sheets_list = sheets_json.get('data', {}).get('sheets', [])
                if sheets_list:
                    actual_sheet_id = sheets_list[0].get('sheet_id')

            if not actual_sheet_id:
                return None

            # 获取表格数据 - 使用 v2 API，range格式为 {sheetId}!A1:Z100
            values_url = f"https://open.feishu.cn/open-apis/sheets/v2/spreadsheets/{spreadsheet_token}/values/{actual_sheet_id}!A1:Z100"
            values_resp = requests.get(values_url, headers=headers, timeout=10)

            if values_resp.status_code != 200:
                return None
            values_json = values_resp.json()
            if values_json.get('code') != 0:
                return None

            # 解析数据为 markdown
            data = values_json.get('data', {})
            value_range = data.get('valueRange', {})
            values = value_range.get('values', [])

            if not values:
                return None

            md_lines = []
            for row in values:
                # 处理每行，将富文本单元格转换为纯文本
                row_cells = []
                for cell in row:
                    if cell is None:
                        row_cells.append('')
                    elif isinstance(cell, str):
                        row_cells.append(cell)
                    elif isinstance(cell, list):
                        # 富文本格式，提取所有 text
                        text_parts = []
                        for segment in cell:
                            if isinstance(segment, dict) and segment.get('type') == 'text':
                                text_parts.append(segment.get('text', ''))
                        row_cells.append(''.join(text_parts))
                    else:
                        row_cells.append(str(cell))
                row_text = ' | '.join(row_cells)
                if row_text.strip():
                    md_lines.append(row_text)

            if md_lines:
                return '## 表格内容\n\n' + '\n'.join(md_lines)
            return None

        except Exception as e:
            return None


class SmartFetcher:
    """智能内容抓取器 - 自动识别链接类型并抓取"""

    def __init__(self):
        self.web_fetcher = WebFetcher()
        self.wechat_fetcher = WechatFetcher()
        self.video_fetcher = VideoFetcher()
        self.feishu_fetcher = FeishuFetcher()
        self.podcast_fetcher = PodcastFetcher()

    def detect_url_type(self, url: str) -> str:
        """检测URL类型"""
        import re

        # 飞书文档
        if 'feishu.cn' in url or 'larksuite.com' in url:
            return 'feishu'

        # 微信公众号
        if 'mp.weixin.qq.com' in url:
            return 'wechat'

        # YouTube
        if re.search(r'(?:youtube\.com|youtu\.be)', url):
            return 'video'

        # Bilibili
        if 'bilibili.com' in url:
            return 'video'

        # 播客RSS
        if '.xml' in url.lower() or 'rss' in url.lower() or 'podcast' in url.lower():
            return 'podcast'

        # 其他网页
        return 'web'

    def fetch(self, url: str) -> dict:
        """智能抓取URL内容"""
        url_type = self.detect_url_type(url)

        if url_type == 'feishu':
            return self.feishu_fetcher.fetch(url)
        elif url_type == 'wechat':
            return self.wechat_fetcher.fetch_article(url)
        elif url_type == 'video':
            return self.video_fetcher.fetch_video_content(url)
        elif url_type == 'podcast':
            return self.podcast_fetcher.fetch_rss(url)
        else:
            return self.web_fetcher.fetch_url(url)


def get_importer(raw_dir: str = None, config: dict = None) -> Importer:
    """获取导入器实例"""
    if raw_dir is None:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        raw_dir = os.path.join(base_dir, 'raw')

    if config is None:
        config_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'config', 'settings.json')
        if os.path.exists(config_path):
            with open(config_path, 'r', encoding='utf-8') as f:
                config = json.load(f)

    return Importer(raw_dir, config)


def get_web_fetcher() -> WebFetcher:
    return WebFetcher()


def get_wechat_fetcher() -> WechatFetcher:
    return WechatFetcher()


def get_video_fetcher() -> VideoFetcher:
    return VideoFetcher()


def get_smart_fetcher() -> SmartFetcher:
    return SmartFetcher()